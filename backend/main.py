import io
import os
import json
from typing import Dict, Tuple
from fastapi import FastAPI, HTTPException, Query
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, validator
from typing import List, Dict, Union, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],  # Adjust this to your React app's URL
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ContentValueType = Union[str, List[str], List[Dict[str, Union[str, int]]]]
ContentValueType = Union[str, List[str], List[List[Union[str, int]]], List[Dict[str, Union[str, int]]]]

class ContentItem(BaseModel):
    type: str
    value: Optional[ContentValueType] = None
    data: Optional[ContentValueType] = None
    headers: Optional[List[str]] = None
    graphType: str = None

class SlideData(BaseModel):
    title: str
    content: List[ContentItem]

class PresentationData(BaseModel):
    slides: List[SlideData]


def text_fits(text_frame, max_size):
    """
    Check if the text fits within the text frame.
    This is an approximation and may not be 100% accurate in all cases.
    """
    if not text_frame.text:
        return True

    text_frame.fit_text(font_family='Arial', max_size=max_size, bold=False, italic=False)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size < 6:  # If font size is too small, consider it doesn't fit
                return False
    return True

def fit_text(text_frame, text_type='normal'):
    """
    Fit text into the text frame by adjusting size and applying wrapping.
    
    :param text_frame: The text frame to fit text into
    :param text_type: Type of text ('title', 'subtitle', or 'normal')
    """
    if text_type == 'title':
        max_size, min_size = 24, 20
    elif text_type == 'subtitle':
        max_size, min_size = 14, 12
    else:  # normal text
        max_size, min_size = 11, 9

    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    # Start with the maximum size
    size = max_size
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(size)

    while not text_fits(text_frame, size) and size > min_size:
        size = max(size * 0.9, min_size)  # Reduce by 10% each iteration
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = size

    # If still doesn't fit, try to wrap (for titles and subtitles)
    if not text_fits(text_frame, size) and text_type in ['title', 'subtitle']:
        words = text_frame.text.split()
        mid = len(words) // 2
        text_frame.text = ' '.join(words[:mid]) + '\n' + ' '.join(words[mid:])
        
        # Readjust size after wrapping
        while not text_fits(text_frame, size) and size > min_size:
            size = max(size * 0.9, min_size)
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = size

    # For normal text, if it still doesn't fit, truncate with ellipsis
    if not text_fits(text_frame, size) and text_type == 'normal':
        while not text_fits(text_frame, size) and len(text_frame.text) > 3:
            text_frame.text = text_frame.text[:-4] + '...'

    return Pt(size)

def create_slide(prs:Presentation, slide_data:SlideData):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Title
    title = slide.shapes.title
    title.text = slide_data.title
    fit_text(title.text_frame, 'title')

    content_top = Inches(1.5)
    for item in slide_data.content:
        if item.type == 'text':
            left, width, height = Inches(0.5), Inches(9), Inches(0.5)
            tf = slide.shapes.add_textbox(left, content_top, width, height).text_frame
            tf.text = item.value
            fit_text(tf, 'normal')
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            content_top += height + Inches(0.1)

        elif item.type == 'bullet':
            shapes = slide.shapes
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = item.value[0]
            for bullet in item.value[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 1
            fit_text(tf, 'normal')
            content_top = Inches(5.5)  # Move to bottom half of slide

        elif item.type == 'graph':
            if item.graphType == 'bar':
                add_chart(slide, item.data, XL_CHART_TYPE.COLUMN_CLUSTERED)
            elif item.graphType == 'pie':
                add_chart(slide, item.data, XL_CHART_TYPE.PIE)
            content_top = Inches(15)  # Move to bottom half of slide
        elif item.type == 'table':
            headers = item.headers
            data = item.data
            rows, cols = len(data) + 1, len(headers)  # +1 for header row
            left, top, width, height = Inches(0.5), content_top, Inches(9), Inches(0.5 * rows)
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Populate the header row
            for col, header in enumerate(headers):
                cell = table.cell(0, col)
                cell.text = str(header)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(10)
            
            # Populate the data rows
            for row, row_data in enumerate(data, start=1):
                for col, cell_data in enumerate(row_data):
                    cell = table.cell(row, col)
                    cell.text = str(cell_data)
                    cell.text_frame.paragraphs[0].font.size = Pt(10)


            content_top += height + Inches(0.1)

def add_chart(slide, data, chart_type):
    chart_data = CategoryChartData()
    chart_data.categories = [item['name'] for item in data]
    chart_data.add_series('Series 1', [item['value'] for item in data])

    # Adjust chart size and position
    x, y = Inches(2.5), Inches(3)
    cx, cy = Inches(5), Inches(2.5)  # Reduced size
    chart = slide.shapes.add_chart(
        chart_type, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(8)  # Adjust legend font size

    # Adjust font sizes and other chart properties
    if chart.chart_title:
        chart.chart_title.text_frame.text = "Chart Title"  # You can customize this
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

    # Adjust series data labels
    for series in chart.series:
        if hasattr(series, 'data_labels'):
            series.data_labels.font.size = Pt(8)
            if chart_type != XL_CHART_TYPE.PIE:
                series.data_labels.position = XL_DATA_LABEL_POSITION.CENTER

    # Adjust axes based on chart type
    if chart_type != XL_CHART_TYPE.PIE:
        # For charts with category and value axes
        if hasattr(chart, 'category_axis') and hasattr(chart.category_axis, 'tick_labels'):
            chart.category_axis.tick_labels.font.size = Pt(8)
        
        if hasattr(chart, 'value_axis') and hasattr(chart.value_axis, 'tick_labels'):
            chart.value_axis.tick_labels.font.size = Pt(8)
            
    else:
        # For pie charts
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.font.size = Pt(8)
        data_labels.position = XL_DATA_LABEL_POSITION.CENTER

    # Additional formatting
    plot = chart.plots[0]
    plot.has_data_labels = True
    if hasattr(plot, 'data_labels'):
        plot.data_labels.font.size = Pt(8)
        plot.data_labels.font.color.rgb = RGBColor(0, 0, 0)  # Black color

# Function 1: Save data to a JSON file
def save_to_json(path: str, filename: str, data: dict):
    full_path = os.path.join(path, filename)
    try:
        os.makedirs(os.path.dirname(full_path), exist_ok=True)
        with open(full_path, 'w') as file:
            json.dump(data, file, indent=2)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error saving file: {str(e)}")

# Function 2: Read data from a JSON file
def read_from_json(path: str, filename: str) -> Tuple[Dict, int, str]:
    full_path = os.path.join(path, filename)
    try:
        with open(full_path, 'r') as file:
            return json.load(file), 200, ""
    except FileNotFoundError:
        return {}, 404, "File not found"
    except json.JSONDecodeError:
        return {}, 400, "Invalid JSON file"
    except Exception as e:
        return {}, 500, f"Error reading file: {str(e)}"

def filter_slides(slides, selected_cards):
    filtered_slides = []

    for slide in slides:
        filtered_content = []
        for content in slide["content"]:
            filtered_data = [
                row for row in content["data"]
                if row[0] in selected_cards
            ]

            if filtered_data:
                filtered_content.append({
                    "headers": content["headers"],
                    "data": filtered_data,
                    "type": content["type"]
                })

        if filtered_content:
            filtered_slides.append({
                "title": slide["title"],
                "content": filtered_content
            })

    return {"slides": filtered_slides}


# Pydantic model for request body
class DataModel(BaseModel):
    data: dict

class GeneratePPTRequest(BaseModel):
    type: str  
    selected_cards: List[str] 

@app.post("/save-business-context")
async def save_business_context(data: DataModel):
    save_to_json("./jsons/input/", "business_context.json", data.data)
    suggestion_output, status_code, error_message = read_from_json("./jsons/output/", "suggestion_output.json")
    if status_code != 200:
        raise HTTPException(status_code=status_code, detail=error_message)
    return suggestion_output

@app.post("/save-suggestions")
async def save_suggestions(data: DataModel):
    save_to_json("./jsons/input/", "suggestion_input.json", data.data)
    
    opportunities, opp_status, opp_error = read_from_json("./jsons/output/", "opportunities.json")
    ai_use, ai_status, ai_error = read_from_json("./jsons/output/", "ai_use.json")
    tech_enablement, tech_status, tech_error = read_from_json("./jsons/output/", "tech_enablement_values.json")
    
    if opp_status != 200:
        raise HTTPException(status_code=opp_status, detail=opp_error)
    if ai_status != 200:
        raise HTTPException(status_code=ai_status, detail=ai_error)
    if tech_status != 200:
        raise HTTPException(status_code=tech_status, detail=tech_error)
    
    return {
        "opportunities": opportunities,
        "ai_responsible_use": ai_use,
        "tech_enablement": tech_enablement
    }

@app.post("/save-tech-reasoning")
async def save_business_context(data: DataModel):
    save_to_json("./jsons/input/", "tech_reasoning.json", data.data)
    suggestion_output, status_code, error_message = read_from_json("./jsons/output/", "result.json")
    if status_code != 200:
        raise HTTPException(status_code=status_code, detail=error_message)
    return suggestion_output

@app.post("/generate-ppt")
async def generate_ppt(request: GeneratePPTRequest):
    # json_data, status_code, error_message = read_from_json("./jsons/output/", "short_table.json")
    if request.type.lower() not in ["short", "comprehensive"]:
        raise HTTPException(status_code=400, detail="Invalid type. Use 'short' or 'comprehensive'.")
    
    json_file = "short_table.json" if request.type == "short" else "table.json"
    filename = "short_presentation.pptx" if request.type == "short" else "presentation.pptx"

    json_data, status_code, error_message = read_from_json("./jsons/output/", json_file)

    # Filter slides
    filtered_json_data = filter_slides(json_data.get("slides", []), request.selected_cards)

    if status_code != 200:
        raise HTTPException(status_code=status_code, detail=error_message)
    
    try:
        presentation_data = PresentationData(**filtered_json_data)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=f"Invalid JSON structure: {str(e)}")

    prs = Presentation()
    for slide_data in presentation_data.slides:
        create_slide(prs, slide_data)

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)

    return StreamingResponse(
        ppt_io,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)